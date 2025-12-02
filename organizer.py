"""Core logic for file processing, classification, and organization.

This module contains functions for:
- Extracting text from various file types (PDF, DOCX, images, etc.).
- Pre-processing images for better OCR results.
- Classifying files based on their content using both local (SBERT) and cloud-based (Google Gemini via Supabase Edge Functions) models.
- Managing a local cache to avoid re-processing files.
- Simulating the file organization process to show users a preview.
"""

import subprocess
import sys
import os
import shutil
import re
import dateparser
import docx
from PIL import Image
import pytesseract
from sentence_transformers import SentenceTransformer, util
import cv2
import numpy as np
from pdf2image import convert_from_path
import fitz
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup
import config
import json
from postgrest.exceptions import APIError
import time
import base64
import hashlib
import requests

# Constants for file extensions, facilitating maintenance
TEXT_BASED_EXTENSIONS = ['.docx', '.pptx', '.xlsx', '.txt', '.html', '.htm']
IMAGE_EXTENSIONS = ['jpg', 'jpeg', 'png', 'bmp', 'gif', 'tiff']
VIDEO_EXTENSIONS = ['mp4', 'mov', 'avi', 'mkv', 'webm']
NATIVE_GEMINI_EXTENSIONS = ['pdf'] + IMAGE_EXTENSIONS + VIDEO_EXTENSIONS


# Patch to ensure that the PyInstaller executable does not open a console window.
if sys.platform == "win32" and getattr(sys, 'frozen', False):
    _original_Popen = subprocess.Popen
    def _patched_popen(*args, **kwargs):
        if 'creationflags' not in kwargs: kwargs['creationflags'] = 0
        kwargs['creationflags'] |= subprocess.CREATE_NO_WINDOW
        return _original_Popen(*args, **kwargs)
    subprocess.Popen = _patched_popen


def get_app_data_path():
    """Creates and returns the path to the DocuSmart data folder in the user's directory.

    This path is OS-dependent to follow platform conventions.

    Returns:
        str: The absolute path to the application's data directory.
    """
    app_name = "DocuSmart"
    if sys.platform == "win32":
        path = os.path.join(os.getenv('APPDATA'), app_name)
    elif sys.platform == "darwin":
        path = os.path.join(os.path.expanduser('~/Library/Application Support'), app_name)
    else:
        path = os.path.join(os.path.expanduser('~/.config'), app_name)
    os.makedirs(path, exist_ok=True)
    return path


def get_file_hash(file_path):
    """Calculates the SHA-256 hash of a file.

    Reads the file in chunks to avoid high memory usage for large files.

    Args:
        file_path (str): The path to the file.

    Returns:
        str or None: The hex digest of the hash, or None if an error occurs.
    """
    sha256_hash = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()
    except Exception as e:
        print(f"AVISO: Não foi possível calcular o hash para {os.path.basename(file_path)}: {e}")
        return None


def load_cache(user_id):
    """Loads the JSON cache file for a specific user.

    The cache stores file hashes and their corresponding classification results.

    Args:
        user_id (str): The unique identifier for the user.

    Returns:
        dict: The loaded cache data, or an empty dictionary if the cache doesn't exist or is invalid.
    """
    cache_path = os.path.join(get_app_data_path(), f"cache_{user_id}.json")
    if os.path.exists(cache_path):
        try:
            with open(cache_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            return {}
    return {}


def save_cache(user_id, cache_data):
    """Saves the cache dictionary to a JSON file for a specific user.

    Args:
        user_id (str): The unique identifier for the user.
        cache_data (dict): The cache dictionary to be saved.
    """
    cache_path = os.path.join(get_app_data_path(), f"cache_{user_id}.json")
    try:
        with open(cache_path, 'w', encoding='utf-8') as f:
            json.dump(cache_data, f, indent=4)
    except IOError as e:
        print(f"ERRO: Falha ao salvar o cache para o usuário {user_id}: {e}")


def get_tesseract_path():
    """Returns the path to the bundled Tesseract executable.

    Detects whether the app is running from source or as a frozen executable (PyInstaller).

    Returns:
        str: The path to the Tesseract executable.
    """
    base_path = sys._MEIPASS if getattr(sys, 'frozen', False) else os.path.dirname(os.path.abspath(__file__))
    if sys.platform == "win32":
        return os.path.join(base_path, 'tesseract', 'tesseract.exe')
    return os.path.join(base_path, 'tesseract', 'tesseract')


def get_poppler_path():
    """Returns the path to the bundled Poppler 'bin' folder.

    Detects whether the app is running from source or as a frozen executable (PyInstaller).
    Required by pdf2image to convert PDFs to images.

    Returns:
        str: The path to the Poppler bin directory.
    """
    base_path = sys._MEIPASS if getattr(sys, 'frozen', False) else os.path.dirname(os.path.abspath(__file__))
    if sys.platform == "win32":
        return os.path.join(base_path, 'poppler-24.08.0', 'Library', 'bin')
    return os.path.join(base_path, 'poppler', 'bin')


def configure_tesseract():
    """Configures environment variables for Tesseract to use local packages.

    Sets the command path for pytesseract and the TESSDATA_PREFIX environment
    variable so Tesseract can find its language data files.
    """
    try:
        tess_path = get_tesseract_path()
        tessdata_path = os.path.join(os.path.dirname(tess_path), 'tessdata')
        pytesseract.pytesseract.tesseract_cmd = tess_path
        os.environ['TESSDATA_PREFIX'] = tessdata_path
    except Exception as e:
        print(f"ERRO CRÍTICO: Falha ao configurar o Tesseract. A função de OCR não funcionará. Detalhes: {e}")


def preprocess_image(pil_image):
    """Applies basic preprocessing to an image to improve OCR results.

    Converts the image to grayscale and applies adaptive thresholding to
    binarize it, which can help Tesseract better distinguish text from the background.

    Args:
        pil_image (PIL.Image.Image): The input image.

    Returns:
        PIL.Image.Image: The preprocessed image.
    """
    try:
        img = np.array(pil_image.convert('RGB')) 
        gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
        thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
        return Image.fromarray(thresh)
    except Exception as e:
        print(f"AVISO: Falha no pré-processamento da imagem: {e}")
        return pil_image 


def invoke_edge_function_manually(function_name, payload, timeout=120):
    """Invokes a Supabase Edge Function via a direct HTTP request with a timeout.

    This is used to call functions that might exceed the default timeout of the
    Supabase Python client library.

    Args:
        function_name (str): The name of the Edge Function.
        payload (dict): The JSON payload to send to the function.
        timeout (int, optional): The request timeout in seconds. Defaults to 120.

    Raises:
        RuntimeError: If the Supabase client is not initialized.
        requests.exceptions.RequestException: For network-related errors.
        requests.exceptions.Timeout: If the request times out.
    """
    if not config.supabase:
        raise RuntimeError("Cliente Supabase não inicializado.")

    url = f"{config.SUPABASE_URL}/functions/v1/{function_name}"
    headers = {
        "Authorization": f"Bearer {config.SUPABASE_KEY}",
        "Content-Type": "application/json"
    }

    try:
        response = requests.post(url, headers=headers, json=payload, timeout=timeout)
        response.raise_for_status()
        return response.json()
    except requests.exceptions.Timeout:
        print(f"ERRO: A requisição para '{function_name}' excedeu o tempo limite de {timeout}s.")
        raise
    except requests.exceptions.RequestException as e:
        print(f"ERRO de requisição ao chamar '{function_name}': {e}")
        raise


def extract_from_pdf(file_path):
    """Extracts text from a PDF file.

    First, it tries to extract text directly. If the extracted text is too short,
    it falls back to using OCR on the rendered pages of the PDF.

    Args:
        file_path (str): The path to the PDF file.

    Returns:
        str: The extracted text content, or an empty string on failure.
    """
    text = ""
    min_text_length_for_ocr_fallback = 100
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text("text")
            
            if len(text.strip()) < min_text_length_for_ocr_fallback and doc.page_count > 0:
                print("  > Texto curto no PDF, tentando OCR como fallback...")
                images = convert_from_path(file_path, poppler_path=get_poppler_path(), dpi=200, last_page=3)
                ocr_text_accumulator = ""
                for i, image in enumerate(images):
                    preprocessed_image = preprocess_image(image)
                    ocr_text_accumulator += pytesseract.image_to_string(preprocessed_image, lang='por+eng')
                
                if len(ocr_text_accumulator.strip()) > len(text.strip()):
                    text = ocr_text_accumulator
    except Exception as e:
        print(f"ERRO ao extrair texto do PDF '{os.path.basename(file_path)}': {e}")
        return ""
    return text.strip()


def extract_from_docx(file_path):
    """Extracts text from a DOCX file.

    Args:
        file_path (str): The path to the DOCX file.

    Returns:
        str: The extracted text content, or an empty string on failure.
    """
    try:
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs]).strip()
    except Exception as e:
        print(f"ERRO ao extrair texto do DOCX '{os.path.basename(file_path)}': {e}")
        return ""


def extract_from_txt(file_path):
    """Extracts text from a plain text file.

    Args:
        file_path (str): The path to the TXT file.

    Returns:
        str: The file's content, or an empty string on failure.
    """
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().strip()
    except Exception as e:
        print(f"ERRO ao extrair texto do TXT '{os.path.basename(file_path)}': {e}")
        return ""


def extract_from_image(file_path):
    """Extracts text from an image file using OCR (Tesseract).

    Args:
        file_path (str): The path to the image file.

    Returns:
        str: The extracted text content, or an empty string on failure.
    """
    try:
        image = Image.open(file_path)
        preprocessed_image = preprocess_image(image)
        return pytesseract.image_to_string(preprocessed_image, lang='por+eng').strip()
    except Exception as e:
        print(f"ERRO ao extrair texto da Imagem '{os.path.basename(file_path)}': {e}")
        return ""


def extract_from_xlsx(file_path):
    """Extracts text from an XLSX (Excel) file.

    Iterates through all cells in all sheets and concatenates their values.

    Args:
        file_path (str): The path to the XLSX file.

    Returns:
        str: The extracted text content, or an empty string on failure.
    """
    text_content = []
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in workbook:
            for row in sheet.iter_rows():
                row_text = [str(cell.value) for cell in row if cell.value is not None]
                if row_text:
                    text_content.append(" | ".join(row_text))
        return "\n".join(text_content).strip()
    except Exception as e:
        print(f"ERRO ao extrair texto do XLSX '{os.path.basename(file_path)}': {e}")
        return ""


def extract_from_pptx(file_path):
    """Extracts text from a PPTX (PowerPoint) file.

    Iterates through all shapes on all slides and extracts their text content.

    Args:
        file_path (str): The path to the PPTX file.

    Returns:
        str: The extracted text content, or an empty string on failure.
    """
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content).strip()
    except Exception as e:
        print(f"ERRO ao extrair texto do PPTX '{os.path.basename(file_path)}': {e}")
        return ""


def extract_from_html(file_path):
    """Extracts text from an HTML file.

    Removes common non-content tags like <script>, <style>, <nav>, etc.,
    before extracting the visible text.

    Args:
        file_path (str): The path to the HTML file.

    Returns:
        str: The extracted text content, or an empty string on failure.
    """
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        for element in soup(["script", "style", "nav", "footer", "aside"]):
            element.extract()
        return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        print(f"ERRO ao extrair texto do HTML '{os.path.basename(file_path)}': {e}")
        return ""


def extract_text_from_file(file_path):
    """Dispatches to the correct text extractor based on the file extension.

    Args:
        file_path (str): The path to the file.

    Returns:
        str: The extracted text, or a message indicating an unsupported format.
    """
    extension = os.path.splitext(file_path)[1].lower()
    
    extraction_map = {
        '.pdf': extract_from_pdf,
        '.docx': extract_from_docx,
        '.txt': extract_from_txt,
        '.xlsx': extract_from_xlsx,
        '.pptx': extract_from_pptx,
        '.html': extract_from_html,
        '.htm': extract_from_html,
    }
    
    if extension in extraction_map:
        return extraction_map[extension](file_path)
    elif extension in [f".{e}" for e in IMAGE_EXTENSIONS]:
        return extract_from_image(file_path)
    else:
        return "Formato de arquivo não suportado."


def extract_dates(text_content):
    """Extracts dates from a block of text using regex and dateparser.

    Args:
        text_content (str): The text to search for dates.

    Returns:
        list[datetime.datetime]: A list of parsed datetime objects found in the text.
    """
    if not text_content: return []
    pattern = r'\b(\d{1,2}[./-]\d{1,2}[./-]\d{2,4})\b'
    found_dates_str = re.findall(pattern, text_content)
    parsed_dates = []
    for date_str_val in found_dates_str:
        try:
            dt = dateparser.parse(date_str_val, languages=['pt', 'en'])
            if dt:
                parsed_dates.append(dt)
        except Exception:
            continue
    return parsed_dates


MODEL_NAME = 'paraphrase-multilingual-mpnet-base-v2'
MODEL_SUBFOLDER = os.path.join('modelos', MODEL_NAME) 
configure_tesseract()

model_sbert = None
try:
    print("Carregando modelo de linguagem local (SBERT)...")
    model_path_sbert = os.path.join(sys._MEIPASS, MODEL_SUBFOLDER) if getattr(sys, 'frozen', False) else MODEL_SUBFOLDER
    if os.path.exists(model_path_sbert):
        model_sbert = SentenceTransformer(model_path_sbert)
        print("Modelo SBERT carregado com sucesso.")
    else:
        print(f"ERRO FATAL: O diretório do modelo SBERT não foi encontrado em '{model_path_sbert}'")
        sys.exit("Falha ao carregar modelo local. O aplicativo não pode continuar.")
except Exception as e:
    print(f"ERRO FATAL ao carregar o modelo de linguagem SBERT: {e}")
    sys.exit("Falha ao carregar modelo local. O aplicativo não pode continuar.")


def classify_content_local(text, categories_embeddings_dict):
    """Classifies text using the local SBERT model via cosine similarity.

    Compares the text's embedding with pre-computed embeddings of category descriptions.

    Args:
        text (str): The text content to classify.
        categories_embeddings_dict (dict): A dictionary mapping category names to their SBERT tensor embeddings.

    Returns:
        tuple[str, float]: A tuple containing the best-matching category name and its confidence score (0.0 to 1.0).
    """
    if not text or len(text.strip()) < 5: return "Outros", 0.0
    text_embedding = model_sbert.encode(text, convert_to_tensor=True)
    max_similarity, best_category = -1.0, "Outros"
    for category_name, description_embedding in categories_embeddings_dict.items():
        if description_embedding is None: continue
        similarity = util.cos_sim(text_embedding, description_embedding).item()
        if similarity > max_similarity:
            max_similarity, best_category = similarity, category_name
    return best_category, (max_similarity + 1) / 2


def classify_by_filename_keywords(filename, categories_dict):
    """Attempts to classify a file based on keywords in its name.

    Uses regex to find common terms associated with predefined categories.

    Args:
        filename (str): The name of the file.
        categories_dict (dict): The dictionary of available categories.

    Returns:
        tuple[str, float]: A tuple containing the matched category and a confidence score (0.98 for category name match, 0.90 for keyword match, 0.0 for no match).
    """
    clean_filename = os.path.splitext(filename)[0].lower().replace("_", " ").replace("-", " ")
    
    keyword_map = {
        "Financeiro": r'\b(extrato|fatura|boleto|conta|holerite|imposto|recibo|nota fiscal|nf-e|nfe|danfe)\b',
        "Pessoal": r'\b(rg|cpf|cnh|passaporte|certidao|eleitor|nascimento|casamento|identidade|titulo eleitoral)\b',
        "Jurídico": r'\b(contrato|peticao|notificacao|judicial|acordo|escritura|procuracao|alvara|sentenca|termo de)\b',
        "Saúde": r'\b(exame|laudo|receita|medico|vacina|consulta|historico medico|atestado|relatorio medico)\b'
    }

    for cat_name in categories_dict.keys():
        pattern_cat = r'\b' + re.escape(cat_name.lower()) + r'\b'
        if re.search(pattern_cat, clean_filename):
            return cat_name, 0.98

    for category, pattern in keyword_map.items():
        if category in categories_dict and re.search(pattern, clean_filename):
            return category, 0.90

    return "Outros", 0.0


def classify_text_via_edge(text_content, categories_dict):
    """Classifies a block of text by calling the 'classify-document-gemini' Edge Function.

    Sends the text content to a Supabase Edge Function for classification by a cloud AI model.

    Args:
        text_content (str): The text to classify.
        categories_dict (dict): The dictionary of available categories to be passed to the AI.

    Returns:
        tuple[str, float]: A tuple containing the predicted category and its confidence score.
    """
    if not config.supabase: return "Outros", 0.0
    if not text_content or len(text_content.strip()) < 10: return "Outros", 0.0

    function_name = "classify-document-gemini"
    payload = {"document_text": text_content, "categories": categories_dict}

    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = invoke_edge_function_manually(function_name, payload, 120)
            if "category" in response:
                return response.get("category", "Outros"), response.get("confidence", 0.3)
            else:
                print(f"ERRO: Resposta da Edge Function (Texto) inesperada: {response}")
                return "Outros", 0.0
        except Exception as e:
            if attempt < max_retries - 1:
                time.sleep(2 * (attempt + 1))
            else:
                print(f"ERRO: Máximo de tentativas atingido para '{function_name}'.")
                raise e


def classify_file_via_edge(file_path, categories_dict):
    """Classifies a file by calling the 'classify-document-file' Edge Function.

    Encodes the file in Base64 and sends it to a Supabase Edge Function for
    classification by a multimodal cloud AI model.

    Args:
        file_path (str): The path to the file to classify.
        categories_dict (dict): The dictionary of available categories to be passed to the AI.

    Returns:
        tuple[str, float]: A tuple containing the predicted category and its confidence score.
    """
    if not config.supabase: return "Outros", 0.0

    function_name = "classify-document-file"
    try:
        with open(file_path, "rb") as f:
            base64_encoded_data = base64.b64encode(f.read()).decode('utf-8')
        
        ext = os.path.splitext(file_path)[1].lower().replace(".", "")
        
        if ext == 'jpg':
            ext = 'jpeg'

        if ext in IMAGE_EXTENSIONS:
            mime_type = f'image/{ext}'
        elif ext in VIDEO_EXTENSIONS:
            mime_type = f'video/{ext}'
        else:
            mime_type = f'application/{ext}'
        
        payload = {
            "file_data_base64": base64_encoded_data,
            "mime_type": mime_type,
            "categories": categories_dict
        }
    except Exception as e:
        print(f"ERRO ao preparar arquivo '{os.path.basename(file_path)}' para upload: {e}")
        return "Outros", 0.0

    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = invoke_edge_function_manually(function_name, payload, 120)
            if "category" in response:
                return response.get("category", "Outros"), response.get("confidence", 0.3)
            else:
                print(f"ERRO: Resposta da Edge Function (Arquivo) inesperada: {response}")
                return "Outros", 0.0
        except Exception as e:
            if attempt < max_retries - 1:
                time.sleep(2 * (attempt + 1))
            else:
                print(f"ERRO: Máximo de tentativas atingido para '{function_name}'.")
                raise e


def simulate_organization(folder_path, categories_dict, progress_callback=None, use_gemini=False, available_credits_for_simulation=0):
    """Simulates the file organization process without moving any files.

    Iterates through files in a given folder, classifies each one using either
    a local model or a cloud AI (Gemini), and returns a structured plan of
    the proposed organization. It uses a cache to speed up re-scans.

    Args:
        folder_path (str): The path to the folder containing files to organize.
        categories_dict (dict): A dictionary of category names to their descriptions.
        progress_callback (function, optional): A callback function to report progress. It receives `current_val`, `total_val`, and `message` arguments.
        use_gemini (bool, optional): If True, attempts to use the Gemini AI for classification.
        available_credits_for_simulation (int, optional): The maximum number of Gemini API calls to make during this simulation.

    Returns:
        tuple: A tuple containing:
            - list: A list of tuples, where each tuple represents a file and its classification result.
            - dict: A dictionary representing the proposed folder structure with files grouped by category.
            - int: The number of Gemini API calls made during the simulation.
    """
    user_id = config.current_user.id if config.current_user else "local_user"
    cache_data = load_cache(user_id)
    cache_was_updated = False

    categories_embeddings_dict = {}
    if not use_gemini or available_credits_for_simulation == 0:
        print("Modo local ativo. Pré-calculando embeddings das categorias...")
        for name, desc in categories_dict.items():
            if desc and name != "Outros (Não processável)":
                categories_embeddings_dict[name] = model_sbert.encode(desc, convert_to_tensor=True)
    
    files_to_organize = []
    organized_structure = {}
    gemini_api_calls_count = 0

    try:
        files_in_folder = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
    except FileNotFoundError:
        return [], {}, 0
        
    total_files = len(files_in_folder)
    if total_files == 0:
        return [], {}, 0

    for i, filename in enumerate(files_in_folder):
        file_path = os.path.join(folder_path, filename)
        extension_with_dot = os.path.splitext(filename)[1].lower()
        extension_no_dot = extension_with_dot.replace(".", "")
        file_hash = get_file_hash(file_path)

        if file_hash and file_hash in cache_data:
            if progress_callback:
                progress_callback(message=f"Verificando cache de '{filename}'...")
            print(f"\nProcessando '{filename}' (Resultado encontrado no cache!)")
            classified_category = cache_data[file_hash]
            classification_method_used = "cache"
        else:
            print(f"\nProcessando '{filename}'...")
            if progress_callback:
                progress_callback(message=f"Analisando '{filename}'...")

            classified_category = "Outros"
            classification_method_used = "não definido"
            gemini_succeeded = False
            
            use_gemini_for_this_file = use_gemini and gemini_api_calls_count < available_credits_for_simulation

            if use_gemini_for_this_file:
                if extension_no_dot in NATIVE_GEMINI_EXTENSIONS:
                    try:
                        print("  > Tentativa 1: IA Gemini (Upload de Arquivo)")
                        category, _ = classify_file_via_edge(file_path, categories_dict)
                        if category != "Outros":
                            classified_category = category
                            classification_method_used = "gemini_file"
                            gemini_succeeded = True
                        else:
                            print("  > Aviso: Upload de arquivo retornou 'Outros'. Tentando fallback.")
                    except Exception as e:
                        print(f"  > Erro no Upload de Arquivo: {e}. Tentando fallback.")
                
                if not gemini_succeeded:
                    try:
                        print("  > Tentativa 2: IA Gemini (Extração de Texto)")
                        text_content = extract_text_from_file(file_path)
                        if text_content and text_content != "Formato de arquivo não suportado.":
                            category, _ = classify_text_via_edge(text_content, categories_dict)
                            if category != "Outros":
                                classified_category = category
                                classification_method_used = "gemini_text"
                                gemini_succeeded = True
                            else:
                                print("  > Aviso: Extração de texto retornou 'Outros'.")
                        else:
                            classified_category = "Outros (Não processável)"
                            classification_method_used = "extracao_falhou"
                            gemini_succeeded = True
                    except Exception as e:
                        print(f"  > Erro na Extração de Texto: {e}.")

                if "gemini" in classification_method_used:
                    gemini_api_calls_count += 1

            if not gemini_succeeded:
                if use_gemini:
                    print("  > Fallback Final: Modelo Local (IA não concluiu)")
                else:
                    print("  > Estratégia: Modelo Local")
                
                kw_category, kw_conf = classify_by_filename_keywords(filename, categories_dict)
                if kw_conf > 0.8:
                    classified_category = kw_category
                    classification_method_used = "local_keyword"
                else:
                    text_content = extract_text_from_file(file_path)
                    if text_content and text_content != "Formato de arquivo não suportado.":
                        classified_category, _ = classify_content_local(text_content, categories_embeddings_dict)
                        classification_method_used = "local_sbert"
                    else:
                        classified_category = "Outros (Não processável)"
                        classification_method_used = "local_nao_processavel"
            
            if file_hash and "gemini" in classification_method_used:
                cache_data[file_hash] = classified_category
                cache_was_updated = True

        date_str = "N/A"
        # try:
        #     # A extração de datas continua sendo um processo separado e informativo
        #     text_for_dates = extract_text_from_file(file_path)
        #     if text_for_dates:
        #         dates_list = extract_dates(text_for_dates)
        #         if dates_list:
        #             date_str = ", ".join(d.strftime('%d/%m/%Y') for d in dates_list)
        # except Exception:
        #     pass

        if classified_category not in categories_dict and classified_category != "Outros (Não processável)":
            classified_category = "Outros"

        if extension_no_dot in IMAGE_EXTENSIONS and classified_category == "Outros":
            if "Imagens" in categories_dict:
                print("  > Ajuste: Imagem classificada como 'Outros'. Redirecionando para 'Imagens'.")
                classified_category = "Imagens"
                classification_method_used += "_img_fallback"
        
        if extension_no_dot in VIDEO_EXTENSIONS and classified_category == "Outros":
            if "Vídeos" in categories_dict:
                print("  > Ajuste: Vídeo classificado como 'Outros'. Redirecionando para 'Vídeos'.")
                classified_category = "Vídeos"
                classification_method_used += "_vid_fallback"

        print(f"  > Resultado Final: Categoria='{classified_category}', Método='{classification_method_used}'")

        files_to_organize.append((filename, classified_category, date_str, classification_method_used))
        if classified_category not in organized_structure:
            organized_structure[classified_category] = []
        organized_structure[classified_category].append(filename)
        
        if progress_callback:
            progress_callback(current_val=i + 1, total_val=total_files)

    if cache_was_updated:
        print("\nSalvando novos resultados no arquivo de cache...")
        save_cache(user_id, cache_data)

    for cat_name_key in categories_dict.keys():
        if cat_name_key not in organized_structure:
            organized_structure[cat_name_key] = []
    
    return files_to_organize, organized_structure, gemini_api_calls_count